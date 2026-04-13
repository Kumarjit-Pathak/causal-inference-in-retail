"""
Generate a polished PowerPoint presentation for the Causal Inference Retail POC.
Reads figures from reports/figures/ and produces reports/Causal_Inference_Tutorial.pptx

Slide deck structure (~25 slides):
  - Title slide
  - Agenda
  - Part 1: Why Causal Inference (3 slides)
  - Part 2: Causal Framework (4 slides)
  - Part 3: Data & Design (3 slides)
  - Part 4: Estimation Methods (4 slides)
  - Part 5: Results (4 slides)
  - Part 6: Robustness (3 slides)
  - Part 7: Recommendations (2 slides)
  - Thank You / Q&A
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ────────────────────────────────────────────────────────────
DARK_BG       = RGBColor(0x1B, 0x1B, 0x2F)   # Deep navy
ACCENT_BLUE   = RGBColor(0x00, 0x96, 0xC7)   # Bright blue
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)   # Orange
ACCENT_GREEN  = RGBColor(0x27, 0xAE, 0x60)   # Green
ACCENT_RED    = RGBColor(0xE7, 0x4C, 0x3C)   # Red
ACCENT_TEAL   = RGBColor(0x00, 0xB4, 0xD8)   # Teal
LIGHT_BG      = RGBColor(0xF8, 0xF9, 0xFA)   # Light grey bg
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GREY          = RGBColor(0x6C, 0x75, 0x7D)
DARK_TEXT      = RGBColor(0x21, 0x25, 0x29)
MID_BLUE      = RGBColor(0x02, 0x3E, 0x8A)

FIGURES = "reports/figures"
OUT = "reports/Causal_Inference_Tutorial.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)  # Widescreen 16:9
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
# Helper Functions
# ═══════════════════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    """Add a plain rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16,
                          color=WHITE, bold=False, line_spacing=1.5, font_name="Calibri",
                          alignment=PP_ALIGN.LEFT):
    """Add text box with multiple lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=16,
                       color=WHITE, bullet_char="\u2022", line_spacing=1.4, font_name="Calibri"):
    """Add bulleted text."""
    lines = [f"{bullet_char}  {b}" for b in bullets]
    return add_multiline_textbox(slide, left, top, width, height, lines,
                                 font_size=font_size, color=color, line_spacing=line_spacing,
                                 font_name=font_name)


def add_image_centered(slide, img_path, top, max_width=Inches(11), max_height=Inches(5)):
    """Add image centered horizontally, scaled to fit."""
    from PIL import Image
    img = Image.open(img_path)
    img_w, img_h = img.size
    aspect = img_w / img_h

    # Scale to fit within bounds
    if max_width / aspect <= max_height:
        width = max_width
        height = int(max_width / aspect)
    else:
        height = max_height
        width = int(max_height * aspect)

    left = int((SLIDE_W - width) / 2)
    slide.shapes.add_picture(img_path, left, top, width, height)


def add_icon_card(slide, left, top, width, height, icon_text, title, desc,
                  bg_color, icon_color=ACCENT_BLUE):
    """Add an infographic card with icon circle, title, and description."""
    # Card background
    card = add_shape(slide, left, top, width, height, bg_color)

    # Icon circle
    circle_size = Inches(0.7)
    circle_left = left + int((width - circle_size) / 2)
    circle_top = top + Inches(0.3)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, circle_size, circle_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = icon_color
    circle.line.fill.background()
    # Icon text
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Segoe UI Emoji"

    # Title
    add_textbox(slide, left + Inches(0.15), top + Inches(1.15), width - Inches(0.3), Inches(0.5),
                title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_textbox(slide, left + Inches(0.15), top + Inches(1.6), width - Inches(0.3), height - Inches(2.0),
                desc, font_size=11, color=RGBColor(0xBB, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)


def section_header_slide(title, subtitle="", slide_number=""):
    """Create a section header slide with gradient-feel background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Accent bar at top
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_BLUE)

    # Left accent line
    add_rect(slide, Inches(1), Inches(2.8), Inches(1.5), Inches(0.06), ACCENT_BLUE)

    # Title
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.2),
                title, font_size=40, color=WHITE, bold=True, font_name="Calibri Light")

    # Subtitle
    if subtitle:
        add_textbox(slide, Inches(1), Inches(4.3), Inches(10), Inches(0.8),
                    subtitle, font_size=20, color=ACCENT_TEAL, font_name="Calibri Light")

    # Slide number
    if slide_number:
        add_textbox(slide, Inches(12), Inches(6.8), Inches(1), Inches(0.5),
                    slide_number, font_size=12, color=GREY, alignment=PP_ALIGN.RIGHT)

    return slide


def content_slide(title, slide_number=""):
    """Create a content slide with dark background and title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

    # Title bar
    add_rect(slide, Inches(0), Inches(0.06), SLIDE_W, Inches(0.95), RGBColor(0x14, 0x14, 0x25))
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                title, font_size=26, color=WHITE, bold=True, font_name="Calibri Light")

    # Slide number
    if slide_number:
        add_textbox(slide, Inches(12), Inches(6.8), Inches(1), Inches(0.5),
                    slide_number, font_size=12, color=GREY, alignment=PP_ALIGN.RIGHT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════════════════════
print("Building slides...")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Decorative shapes
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.1), ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), ACCENT_BLUE)

# Left decorative bar
add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.08), Inches(3.5), ACCENT_TEAL)

# Title
add_textbox(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(1.0),
            "Causal Inference for Retail Promotions",
            font_size=44, color=WHITE, bold=True, font_name="Calibri Light")

# Subtitle
add_textbox(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(0.8),
            "A Complete Technical Tutorial",
            font_size=28, color=ACCENT_TEAL, font_name="Calibri Light")

# Details
add_multiline_textbox(slide, Inches(1.2), Inches(4.3), Inches(8), Inches(1.5),
    ["4 Indian Cities  |  100 SKUs  |  3 Years of Weekly Data",
     "DoWhy + EconML  |  DAGs + Double Machine Learning",
     "From Correlation to Causation to Strategy"],
    font_size=16, color=GREY, line_spacing=1.6)

# Bottom tag
add_textbox(slide, Inches(1.2), Inches(6.5), Inches(5), Inches(0.5),
            "Target Audience: Data Scientists & ML Engineers",
            font_size=14, color=RGBColor(0x88, 0x88, 0x99))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Agenda", "2")

agenda_items = [
    ("1", "Why Causal Inference?", "The correlation trap and why standard ML fails", ACCENT_BLUE),
    ("2", "The Causal Framework", "DAGs, Backdoor Criterion, d-separation", ACCENT_TEAL),
    ("3", "Data & Experimental Design", "Synthetic data with known causal structure", ACCENT_ORANGE),
    ("4", "Estimation Methodology", "DoWhy, DML, CausalForest", ACCENT_GREEN),
    ("5", "Results & Findings", "City-level uplift, CATE, winning combos", RGBColor(0x9B, 0x59, 0xB6)),
    ("6", "Robustness & Sensitivity", "Refutation, confounder omission, stability", ACCENT_RED),
    ("7", "Strategic Recommendations", "Actionable playbook by city", RGBColor(0xF3, 0x9C, 0x12)),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    row = i
    y = Inches(1.3) + row * Inches(0.82)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"

    # Title
    add_textbox(slide, Inches(1.8), y - Inches(0.02), Inches(4), Inches(0.4),
                title, font_size=18, color=WHITE, bold=True)
    # Description
    add_textbox(slide, Inches(1.8), y + Inches(0.32), Inches(6), Inches(0.35),
                desc, font_size=13, color=GREY)

    # Connecting line
    if i < len(agenda_items) - 1:
        add_rect(slide, Inches(1.25), y + Inches(0.55), Inches(0.03), Inches(0.27), RGBColor(0x33, 0x33, 0x55))

# Right side: key stats
stats_x = Inches(8.5)
add_shape(slide, stats_x, Inches(1.3), Inches(4.0), Inches(5.5), RGBColor(0x14, 0x14, 0x25),
          border_color=RGBColor(0x33, 0x33, 0x55), border_width=Pt(1))

add_textbox(slide, stats_x + Inches(0.3), Inches(1.5), Inches(3.5), Inches(0.5),
            "Project At a Glance", font_size=18, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

stats = [
    ("62,400", "Observations"),
    ("7", "Treatment Levers"),
    ("4", "Indian Cities"),
    ("156", "Weeks (3 Years)"),
    ("100", "SKUs Analyzed"),
    ("6", "Robustness Tests"),
]
for i, (val, label) in enumerate(stats):
    sy = Inches(2.2) + i * Inches(0.72)
    add_textbox(slide, stats_x + Inches(0.5), sy, Inches(1.5), Inches(0.4),
                val, font_size=28, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, stats_x + Inches(2.1), sy + Inches(0.05), Inches(2.0), Inches(0.4),
                label, font_size=14, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# PART 1: WHY CAUSAL INFERENCE
# ═══════════════════════════════════════════════════════════════════════════════
section_header_slide("Part 1", "Why Causal Inference?", "3")

# Slide: The Correlation Trap
slide = content_slide("The Correlation Trap", "4")

# Three cards showing bias types
cards = [
    ("!", "Confounding\nBias", "High-equity brands get more promos AND sell more.\nNaive models attribute the brand effect to the promotion.",
     ACCENT_RED),
    ("?", "Simpson's\nParadox", "A discount looks effective overall but harmful in every city\n-- because discounts cluster in high-sales cities.",
     ACCENT_ORANGE),
    ("X", "Selection\nBias", "Stores with displays are already top performers.\nThe 'display effect' is partly pre-existing advantage.",
     RGBColor(0x9B, 0x59, 0xB6)),
]

for i, (icon, title, desc, color) in enumerate(cards):
    x = Inches(0.8) + i * Inches(4.1)
    card = add_shape(slide, x, Inches(1.5), Inches(3.8), Inches(3.5),
                     RGBColor(0x14, 0x14, 0x25), border_color=color, border_width=Pt(2))

    # Icon
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.4), Inches(1.8), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.2), Inches(2.9), Inches(3.4), Inches(0.7),
                title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(3.6), Inches(3.4), Inches(1.2),
                desc, font_size=11, color=RGBColor(0xBB, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

# Bottom message
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(11), Inches(0.8),
            "Causal inference solves this by formally identifying and adjusting for confounding paths.",
            font_size=18, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)


# Slide: Naive vs Causal comparison
slide = content_slide("The Bias We're Correcting", "5")

# Comparison table
treatments = [
    ("discount_depth", "520", "468", "-10.0%"),
    ("is_instore_display", "195", "164", "-15.9%"),
    ("local_channel_promo", "98", "79", "-19.4%"),
    ("sms_blast_active", "110", "83", "-24.5%"),
    ("loyalty_topup_discount", "180", "151", "-16.1%"),
    ("special_coupon_usage", "260", "215", "-17.3%"),
    ("is_2x_points_active", "22", "15", "-31.8%"),
]

# Table header
headers = ["Treatment Lever", "Naive Effect", "Causal ATE", "Bias"]
header_colors = [ACCENT_BLUE, ACCENT_ORANGE, ACCENT_GREEN, ACCENT_RED]
table_left = Inches(1.5)
table_top = Inches(1.5)
col_widths = [Inches(3.5), Inches(2.2), Inches(2.2), Inches(2.2)]

for j, (header, hcolor) in enumerate(zip(headers, header_colors)):
    x = table_left + sum(cw for cw in col_widths[:j])
    rect = add_rect(slide, x, table_top, col_widths[j], Inches(0.55), RGBColor(0x14, 0x14, 0x25))
    add_textbox(slide, x, table_top + Inches(0.05), col_widths[j], Inches(0.45),
                header, font_size=14, color=hcolor, bold=True, alignment=PP_ALIGN.CENTER)

# Table rows
for i, (name, naive, causal, bias) in enumerate(treatments):
    y = table_top + Inches(0.55) + i * Inches(0.6)
    bg = RGBColor(0x1E, 0x1E, 0x35) if i % 2 == 0 else RGBColor(0x17, 0x17, 0x2A)
    for j, val in enumerate([name, naive, causal, bias]):
        x = table_left + sum(cw for cw in col_widths[:j])
        add_rect(slide, x, y, col_widths[j], Inches(0.55), bg)
        clr = WHITE if j == 0 else (ACCENT_ORANGE if j == 1 else (ACCENT_GREEN if j == 2 else ACCENT_RED))
        fs = 13 if j == 0 else 15
        add_textbox(slide, x, y + Inches(0.07), col_widths[j], Inches(0.4),
                    val, font_size=fs, color=clr, bold=(j > 0), alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# Insight callout
add_shape(slide, Inches(1.5), Inches(5.9), Inches(10.1), Inches(0.9), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_RED, border_width=Pt(2))
add_textbox(slide, Inches(1.8), Inches(6.0), Inches(9.5), Inches(0.7),
            "Every naive estimate is inflated by confounding. Smaller true effects have larger proportional bias (2x Points: 32% overstated).",
            font_size=14, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# PART 2: CAUSAL FRAMEWORK
# ═══════════════════════════════════════════════════════════════════════════════
section_header_slide("Part 2", "The Causal Framework", "6")

# Slide: Potential Outcomes
slide = content_slide("Potential Outcomes Framework (Neyman-Rubin)", "7")

add_multiline_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5), [
    "For each store-week observation i :",
    "",
    "  Y_i(1)  =  Sales if treated (e.g., SMS sent)",
    "  Y_i(0)  =  Sales if NOT treated",
    "",
    "Individual Treatment Effect:",
    "  ITE_i  =  Y_i(1) - Y_i(0)",
    "",
    "Average Treatment Effect:",
    "  ATE  =  E[Y(1)] - E[Y(0)]",
], font_size=15, color=WHITE, line_spacing=1.3, font_name="Consolas")

# Right side: fundamental problem
add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(2.5), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_ORANGE, border_width=Pt(2))
add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.0), Inches(0.5),
            "The Fundamental Problem", font_size=18, color=ACCENT_ORANGE, bold=True)
add_multiline_textbox(slide, Inches(7.3), Inches(2.1), Inches(5.0), Inches(1.5), [
    "We NEVER observe both Y(1) and Y(0)",
    "for the same unit at the same time.",
    "",
    "One potential outcome is always",
    "counterfactual -- we must estimate it.",
], font_size=14, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.4)

# Bottom: solution
add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_GREEN, border_width=Pt(2))
add_textbox(slide, Inches(1.1), Inches(4.4), Inches(11.0), Inches(0.5),
            "Solution: Conditional Independence (Ignorability)", font_size=18, color=ACCENT_GREEN, bold=True)
add_multiline_textbox(slide, Inches(1.1), Inches(5.0), Inches(11.0), Inches(1.5), [
    "If we condition on the right set of confounders W, the treatment assignment becomes 'as-if random':",
    "    Y(0), Y(1)  independent of  T  |  W",
    "",
    "This allows us to estimate ATE = E[Y | T=1, W] - E[Y | T=0, W], averaged over W.",
    "The DAG tells us exactly which variables belong in W."
], font_size=14, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.3, font_name="Calibri")


# Slide: DAG
slide = content_slide("The Causal DAG", "8")
add_image_centered(slide, f"{FIGURES}/dag_publication.png", Inches(1.2), max_width=Inches(11.5), max_height=Inches(5.5))


# Slide: Backdoor Criterion
slide = content_slide("Backdoor Criterion & d-Separation", "9")

# Three d-separation patterns
patterns = [
    ("Chain", "A -> B -> C", "Conditioning on B\nBLOCKS the path", ACCENT_BLUE),
    ("Fork", "A <- B -> C", "Conditioning on B\nBLOCKS the path", ACCENT_GREEN),
    ("Collider", "A -> B <- C", "Conditioning on B\nOPENS the path!", ACCENT_RED),
]

for i, (name, graph, desc, color) in enumerate(patterns):
    x = Inches(0.6) + i * Inches(4.2)
    card = add_shape(slide, x, Inches(1.3), Inches(3.8), Inches(2.5),
                     RGBColor(0x14, 0x14, 0x25), border_color=color, border_width=Pt(2))
    add_textbox(slide, x, Inches(1.4), Inches(3.8), Inches(0.5),
                name, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(1.95), Inches(3.8), Inches(0.5),
                graph, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_textbox(slide, x, Inches(2.6), Inches(3.8), Inches(0.8),
                desc, font_size=13, color=RGBColor(0xBB, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

# Backdoor criterion box
add_shape(slide, Inches(0.6), Inches(4.2), Inches(12.0), Inches(2.8), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_TEAL, border_width=Pt(2))
add_textbox(slide, Inches(0.9), Inches(4.4), Inches(11.0), Inches(0.5),
            "Pearl's Backdoor Criterion", font_size=20, color=ACCENT_TEAL, bold=True)
add_multiline_textbox(slide, Inches(0.9), Inches(5.0), Inches(5.5), Inches(1.8), [
    "A set Z satisfies the backdoor criterion if:",
    "  1. No node in Z is a descendant of T",
    "  2. Z blocks every backdoor path between T and Y",
    "",
    "Conditioning on Z gives unbiased causal estimates.",
], font_size=14, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.3)

add_multiline_textbox(slide, Inches(7.0), Inches(5.0), Inches(5.5), Inches(1.8), [
    "Our adjustment set:",
    "  Z = { brand_equity,",
    "         competitor_price_index,",
    "         seasonality_multiplier,",
    "         is_festival_week,",
    "         base_price, city_id }",
], font_size=14, color=ACCENT_GREEN, line_spacing=1.2, font_name="Consolas")


# ═══════════════════════════════════════════════════════════════════════════════
# PART 3: DATA & DESIGN
# ═══════════════════════════════════════════════════════════════════════════════
section_header_slide("Part 3", "Data & Experimental Design", "10")

# Slide: Data overview
slide = content_slide("Synthetic Data with Known Causal Structure", "11")

# Why synthetic data box
add_shape(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(1.8), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_BLUE, border_width=Pt(2))
add_textbox(slide, Inches(0.9), Inches(1.35), Inches(5.3), Inches(0.4),
            "Why Synthetic Data?", font_size=16, color=ACCENT_BLUE, bold=True)
add_multiline_textbox(slide, Inches(0.9), Inches(1.8), Inches(5.3), Inches(1.0), [
    "We know the TRUE causal effects by construction.",
    "This lets us validate that our causal methods",
    "actually recover the correct treatment effects.",
], font_size=13, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.3)

# Data specs
add_shape(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(1.8), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_GREEN, border_width=Pt(2))
add_textbox(slide, Inches(7.1), Inches(1.35), Inches(5.3), Inches(0.4),
            "Dataset Specifications", font_size=16, color=ACCENT_GREEN, bold=True)
add_multiline_textbox(slide, Inches(7.1), Inches(1.8), Inches(5.3), Inches(1.0), [
    "62,400 rows = 100 SKUs x 156 weeks x 4 cities",
    "24 columns: 7 treatments, 6 confounders, 3 outcomes",
    "Seed: np.random.seed(42) for full reproducibility",
], font_size=13, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.3)

# Non-linear interactions
add_shape(slide, Inches(0.6), Inches(3.4), Inches(12.0), Inches(3.5), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_ORANGE, border_width=Pt(2))
add_textbox(slide, Inches(0.9), Inches(3.55), Inches(11.0), Inches(0.4),
            "Embedded Non-Linear Interactions (Key Design Choices)", font_size=16, color=ACCENT_ORANGE, bold=True)

interactions = [
    ("SMS + Display Synergy", "synergy = 90 * sms_blast * display", "+90 units when both active (36% above additive)", ACCENT_BLUE),
    ("Discount Saturation", "effect = 20*d - 0.3*d^2", "Quadratic: diminishing returns beyond 20%", ACCENT_RED),
    ("Loyalty Interaction", "boost = 30 * is_2x * coupon", "2x Points amplifies coupon redemption effect", ACCENT_GREEN),
    ("Confounded Assignment", "discount ~ brand_equity + festival", "High-equity brands get deeper discounts", ACCENT_ORANGE),
]

for i, (name, code, desc, color) in enumerate(interactions):
    y = Inches(4.1) + i * Inches(0.65)
    # Color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.1), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_textbox(slide, Inches(1.6), y - Inches(0.02), Inches(2.8), Inches(0.4),
                name, font_size=13, color=WHITE, bold=True)
    add_textbox(slide, Inches(4.5), y - Inches(0.02), Inches(3.5), Inches(0.4),
                code, font_size=12, color=ACCENT_TEAL, font_name="Consolas")
    add_textbox(slide, Inches(8.2), y - Inches(0.02), Inches(4.2), Inches(0.4),
                desc, font_size=12, color=GREY)


# Slide: Festival Seasonality
slide = content_slide("Indian Festival Seasonality", "12")

festivals = [
    ("Diwali", "Oct-Nov (wk 44-45)", "Mumbai, Delhi", "1.0", ACCENT_ORANGE),
    ("Holi", "March (wk 11-12)", "Delhi, Mumbai", "0.9", RGBColor(0xE9, 0x1E, 0x63)),
    ("Eid", "Varies by year", "Hyderabad", "1.0", ACCENT_GREEN),
    ("Pongal", "January (wk 2-3)", "Bengaluru", "1.0", ACCENT_BLUE),
    ("Christmas", "December (wk 52)", "All cities", "0.7", ACCENT_RED),
    ("Ganesh Chaturthi", "Aug-Sep", "Mumbai", "1.0", ACCENT_TEAL),
    ("Navratri", "October", "All cities", "0.8", RGBColor(0x9B, 0x59, 0xB6)),
]

# Header
fheaders = ["Festival", "Timing", "Primary Cities", "Weight"]
fx = [Inches(1.0), Inches(3.8), Inches(6.5), Inches(9.8)]
fw = [Inches(2.8), Inches(2.7), Inches(3.3), Inches(2.0)]
for j, (h, hx, hw) in enumerate(zip(fheaders, fx, fw)):
    add_rect(slide, hx, Inches(1.3), hw, Inches(0.5), ACCENT_BLUE)
    add_textbox(slide, hx, Inches(1.33), hw, Inches(0.45),
                h, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (name, timing, cities, weight, color) in enumerate(festivals):
    y = Inches(1.8) + i * Inches(0.6)
    bg = RGBColor(0x1E, 0x1E, 0x35) if i % 2 == 0 else RGBColor(0x17, 0x17, 0x2A)
    for j, (val, vx, vw) in enumerate(zip([name, timing, cities, weight], fx, fw)):
        add_rect(slide, vx, y, vw, Inches(0.55), bg)
        clr = color if j == 0 else WHITE
        add_textbox(slide, vx, y + Inches(0.07), vw, Inches(0.4),
                    val, font_size=13, color=clr, bold=(j == 0), alignment=PP_ALIGN.CENTER)

# Insight box
add_shape(slide, Inches(1.0), Inches(6.0), Inches(10.8), Inches(0.8), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_ORANGE, border_width=Pt(1))
add_textbox(slide, Inches(1.3), Inches(6.1), Inches(10.2), Inches(0.6),
            "Festival weeks amplify BOTH treatment propensity AND base sales, creating seasonal confounding that must be controlled for.",
            font_size=14, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# PART 4: ESTIMATION METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════════
section_header_slide("Part 4", "Estimation Methodology", "13")

# Slide: DoWhy Pipeline
slide = content_slide("DoWhy: The 4-Step Causal Pipeline", "14")

steps = [
    ("1", "MODEL", "Define the causal DAG\n(nodes + directed edges)", ACCENT_BLUE),
    ("2", "IDENTIFY", "Find estimand via\nBackdoor Criterion", ACCENT_TEAL),
    ("3", "ESTIMATE", "Compute ATE/CATE\nusing regression/DML", ACCENT_GREEN),
    ("4", "REFUTE", "Validate with placebo,\nrandom CC, subset tests", ACCENT_ORANGE),
]

for i, (num, name, desc, color) in enumerate(steps):
    x = Inches(0.5) + i * Inches(3.2)
    card = add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(3.0),
                     RGBColor(0x14, 0x14, 0x25), border_color=color, border_width=Pt(2))

    # Step number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.0), Inches(1.8), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, x, Inches(2.8), Inches(2.9), Inches(0.5),
                name, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(3.3), Inches(2.5), Inches(1.0),
                desc, font_size=13, color=RGBColor(0xBB, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < 3:
        arrow_x = x + Inches(3.0)
        add_textbox(slide, arrow_x, Inches(2.6), Inches(0.3), Inches(0.5),
                    ">", font_size=28, color=GREY, bold=True, alignment=PP_ALIGN.CENTER)

# Code snippet
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), RGBColor(0x0D, 0x0D, 0x1A))
add_multiline_textbox(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.8), [
    'model = dowhy.CausalModel(data=df, treatment="discount_depth", outcome="sales_volume", graph=gml)',
    'estimand = model.identify_effect()              # Step 2: Backdoor identification',
    'estimate = model.estimate_effect(estimand,       # Step 3: Linear regression estimate',
    '               method_name="backdoor.linear_regression")',
    'refutation = model.refute_estimate(estimand, estimate,  # Step 4: Placebo test',
    '               method_name="placebo_treatment_refuter")',
], font_size=11, color=ACCENT_TEAL, line_spacing=1.2, font_name="Consolas")


# Slide: DML Deep Dive
slide = content_slide("Double Machine Learning (Chernozhukov et al., 2018)", "15")

# Three-step visual
dml_steps = [
    ("Step 1", "Residualize Y", "Y_res = Y - E[Y|W]\n\nRemove confounder\neffects from outcome\nusing ML (GBM)", ACCENT_BLUE),
    ("Step 2", "Residualize T", "T_res = T - E[T|W]\n\nRemove confounder\neffects from treatment\nusing ML (GBM)", ACCENT_ORANGE),
    ("Step 3", "Causal Regression", "Y_res = theta * T_res\n\ntheta is the CAUSAL\neffect, immune to\nregularization bias", ACCENT_GREEN),
]

for i, (step, title, desc, color) in enumerate(dml_steps):
    x = Inches(0.5) + i * Inches(4.2)
    card = add_shape(slide, x, Inches(1.3), Inches(3.8), Inches(3.2),
                     RGBColor(0x14, 0x14, 0x25), border_color=color, border_width=Pt(2))
    add_textbox(slide, x, Inches(1.4), Inches(3.8), Inches(0.4),
                step, font_size=12, color=GREY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(1.8), Inches(3.8), Inches(0.5),
                title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_textbox(slide, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(1.8),
                          desc.split("\n"), font_size=13, color=RGBColor(0xCC, 0xCC, 0xDD),
                          line_spacing=1.2, alignment=PP_ALIGN.CENTER, font_name="Consolas")

    if i < 2:
        add_textbox(slide, x + Inches(3.85), Inches(2.5), Inches(0.4), Inches(0.5),
                    ">", font_size=28, color=GREY, bold=True)

# Why DML box
add_shape(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.2), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_TEAL, border_width=Pt(2))
add_textbox(slide, Inches(0.8), Inches(5.05), Inches(5), Inches(0.4),
            "Why DML over Plain Regression?", font_size=16, color=ACCENT_TEAL, bold=True)

dml_advantages = [
    "Handles non-linear confounding (ML first stage handles complexity)",
    "Valid confidence intervals via Neyman-orthogonal construction",
    "Cross-fitting prevents overfitting bias in high dimensions",
    "Automatic heterogeneous effects (CATE) via effect modifiers X",
]
add_bullet_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
                   dml_advantages, font_size=13, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.3)


# Slide: CausalForest
slide = content_slide("CausalForest vs LinearDML", "16")

# Comparison cards
comparisons = [
    ("LinearDML", [
        "Effect is LINEAR in X",
        "theta(x) = beta_0 + beta_1 * x",
        "Fast, interpretable",
        "Best for: known effect modifiers",
        "Our use: multi-treatment CATE",
    ], ACCENT_BLUE),
    ("CausalForestDML", [
        "Effect is NON-LINEAR in X",
        "theta(x) = learned tree ensemble",
        "Flexible, less interpretable",
        "Best for: discovering heterogeneity",
        "Our use: binary treatment CATE",
    ], ACCENT_GREEN),
]

for i, (name, bullets, color) in enumerate(comparisons):
    x = Inches(0.6) + i * Inches(6.3)
    card = add_shape(slide, x, Inches(1.3), Inches(5.8), Inches(4.0),
                     RGBColor(0x14, 0x14, 0x25), border_color=color, border_width=Pt(2))
    add_textbox(slide, x, Inches(1.5), Inches(5.8), Inches(0.5),
                name, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_textbox(slide, x + Inches(0.3), Inches(2.2), Inches(5.2), Inches(2.5),
                       bullets, font_size=14, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.5)

# Technical gotchas
add_shape(slide, Inches(0.6), Inches(5.7), Inches(12.0), Inches(1.2), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_RED, border_width=Pt(1))
add_textbox(slide, Inches(0.9), Inches(5.8), Inches(2.0), Inches(0.4),
            "Gotchas:", font_size=14, color=ACCENT_RED, bold=True)
add_multiline_textbox(slide, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.7), [
    "model_t MUST be a regressor (not classifier) -- even for binary T.  |  "
    "n_estimators must be divisible by subforest_size (4).  |  "
    "Multi-treatment: wrap model_t in MultiOutputRegressor()."
], font_size=12, color=RGBColor(0xBB, 0xBB, 0xCC), line_spacing=1.2)


# ═══════════════════════════════════════════════════════════════════════════════
# PART 5: RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
section_header_slide("Part 5", "Results & Findings", "17")

# Slide: City Uplift Bars
slide = content_slide("City-Level Causal Uplift (ATE)", "18")
add_image_centered(slide, f"{FIGURES}/city_uplift_bars.png", Inches(1.2), max_width=Inches(12.0), max_height=Inches(5.5))

# Slide: CATE Distributions
slide = content_slide("Treatment Effect Heterogeneity (CATE by City)", "19")
add_image_centered(slide, f"{FIGURES}/cate_distributions.png", Inches(1.2), max_width=Inches(12.0), max_height=Inches(5.5))

# Slide: Winning Combo
slide = content_slide("Winning Combination Matrix", "20")
add_image_centered(slide, f"{FIGURES}/winning_combo_heatmap.png", Inches(1.2), max_width=Inches(12.0), max_height=Inches(4.5))

# Insight at bottom
add_shape(slide, Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.9), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_ORANGE, border_width=Pt(1))
add_textbox(slide, Inches(0.9), Inches(6.1), Inches(11.4), Inches(0.7),
            "Stars mark the top-performing lever per city. Discount Depth dominates volume, but erodes revenue. "
            "SMS + Display synergy yields +90 bonus units (36% above additive) -- always pair them together.",
            font_size=13, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)


# Slide: Key insight - revenue vs volume
slide = content_slide("The Volume vs. Revenue Trade-Off", "21")

# Left: volume winners
add_shape(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.2), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_GREEN, border_width=Pt(2))
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.5),
            "Volume Drivers", font_size=22, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
vol_items = [
    ("1. Coupon Usage", "+470 units/week", "Highest & most consistent"),
    ("2. Discount Depth", "+468 units/pct point", "BUT destroys revenue!"),
    ("3. In-Store Display", "+164 units/week", "Also boosts revenue"),
    ("4. Loyalty Top-Up", "+151 units/week", "Moderate, consistent"),
    ("5. SMS Blast", "+83 units/week", "Best paired with Display"),
]
for i, (name, effect, note) in enumerate(vol_items):
    y = Inches(2.2) + i * Inches(0.78)
    add_textbox(slide, Inches(0.9), y, Inches(2.8), Inches(0.35),
                name, font_size=14, color=WHITE, bold=True)
    add_textbox(slide, Inches(3.7), y, Inches(1.6), Inches(0.35),
                effect, font_size=13, color=ACCENT_GREEN, bold=True)
    clr = ACCENT_RED if "destroys" in note else GREY
    add_textbox(slide, Inches(3.7), y + Inches(0.3), Inches(2.5), Inches(0.35),
                note, font_size=11, color=clr)

# Right: revenue impact
add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.2), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_RED, border_width=Pt(2))
add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5),
            "Revenue Impact Warning", font_size=22, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)

add_multiline_textbox(slide, Inches(7.2), Inches(2.3), Inches(5.0), Inches(3.0), [
    "Discount Depth:",
    "  Volume ATE:  +468 units",
    "  Revenue ATE: -70,000 INR",
    "",
    "The margin erosion FAR outweighs",
    "the volume gain at scale.",
    "",
    "In contrast, In-Store Display:",
    "  Volume ATE:  +164 units",
    "  Revenue ATE: +20,000 INR",
    "",
    "Display is the most BALANCED lever:",
    "drives both volume AND revenue.",
], font_size=14, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.2)


# ═══════════════════════════════════════════════════════════════════════════════
# PART 6: ROBUSTNESS
# ═══════════════════════════════════════════════════════════════════════════════
section_header_slide("Part 6", "Robustness & Sensitivity Analysis", "22")

# Slide: Robustness Scorecard
slide = content_slide("Robustness Scorecard", "23")
add_image_centered(slide, f"{FIGURES}/robustness_scorecard.png", Inches(1.2), max_width=Inches(12.0), max_height=Inches(3.5))

# Explanation below
add_shape(slide, Inches(0.6), Inches(5.0), Inches(3.8), Inches(2.0), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_BLUE, border_width=Pt(1))
add_textbox(slide, Inches(0.8), Inches(5.1), Inches(3.4), Inches(0.3),
            "Placebo Test", font_size=14, color=ACCENT_BLUE, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.45), Inches(3.4), Inches(1.3),
            "Randomly permute treatment. If effect drops to ~0, the original was real.",
            font_size=12, color=RGBColor(0xBB, 0xBB, 0xCC))

add_shape(slide, Inches(4.7), Inches(5.0), Inches(3.8), Inches(2.0), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_GREEN, border_width=Pt(1))
add_textbox(slide, Inches(4.9), Inches(5.1), Inches(3.4), Inches(0.3),
            "Random Common Cause", font_size=14, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(4.9), Inches(5.45), Inches(3.4), Inches(1.3),
            "Add random noise confounder. If ATE barely changes (<5%), estimate is robust.",
            font_size=12, color=RGBColor(0xBB, 0xBB, 0xCC))

add_shape(slide, Inches(8.8), Inches(5.0), Inches(3.8), Inches(2.0), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_ORANGE, border_width=Pt(1))
add_textbox(slide, Inches(9.0), Inches(5.1), Inches(3.4), Inches(0.3),
            "Data Subset Test", font_size=14, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(9.0), Inches(5.45), Inches(3.4), Inches(1.3),
            "Re-estimate on 70% subsample. If ATE is stable (<5% change), not driven by outliers.",
            font_size=12, color=RGBColor(0xBB, 0xBB, 0xCC))


# Slide: Sensitivity Tornado
slide = content_slide("Confounder Sensitivity Analysis", "24")
add_image_centered(slide, f"{FIGURES}/sensitivity_tornado.png", Inches(1.2), max_width=Inches(12.0), max_height=Inches(4.0))

add_shape(slide, Inches(0.6), Inches(5.6), Inches(12.0), Inches(1.3), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_RED, border_width=Pt(1))
add_multiline_textbox(slide, Inches(0.9), Inches(5.7), Inches(11.4), Inches(1.1), [
    "Brand equity is the most critical confounder -- omitting it shifts Discount Depth ATE by 35%.",
    "In-Store Display and Coupon Usage are robust to confounder omission (<3% shift).",
    "Any production model MUST include brand equity as a control variable.",
], font_size=13, color=ACCENT_RED, line_spacing=1.4)


# Slide: CATE Stability
slide = content_slide("CATE Cross-Split Stability", "25")

# Stability table
stab_data = [
    ("Mumbai", "163.0", "165.5", "160.5", "1.6%"),
    ("Bengaluru", "166.4", "163.0", "169.8", "2.0%"),
    ("Delhi", "165.5", "166.8", "164.3", "0.8%"),
    ("Hyderabad", "163.3", "162.5", "164.2", "0.5%"),
]
stab_headers = ["City", "Split 1", "Split 2", "Split 3", "CV%"]
stab_colors = [ACCENT_BLUE, WHITE, WHITE, WHITE, ACCENT_GREEN]
sx = Inches(2.5)
scols = [Inches(2.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(1.5)]

for j, (h, sc) in enumerate(zip(stab_headers, stab_colors)):
    cx = sx + sum(scols[:j])
    add_rect(slide, cx, Inches(1.5), scols[j], Inches(0.55), ACCENT_BLUE)
    add_textbox(slide, cx, Inches(1.55), scols[j], Inches(0.45),
                h, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(stab_data):
    y = Inches(2.05) + i * Inches(0.6)
    bg = RGBColor(0x1E, 0x1E, 0x35) if i % 2 == 0 else RGBColor(0x17, 0x17, 0x2A)
    for j, val in enumerate(row):
        cx = sx + sum(scols[:j])
        add_rect(slide, cx, y, scols[j], Inches(0.55), bg)
        clr = stab_colors[j]
        add_textbox(slide, cx, y + Inches(0.08), scols[j], Inches(0.4),
                    val, font_size=14, color=clr, bold=(j == 4), alignment=PP_ALIGN.CENTER)

# Conclusion box
add_shape(slide, Inches(2.0), Inches(4.6), Inches(9.3), Inches(1.2), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_GREEN, border_width=Pt(2))
add_textbox(slide, Inches(2.3), Inches(4.75), Inches(8.7), Inches(0.4),
            "All CV% values are well below 20%", font_size=18, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.3), Inches(5.2), Inches(8.7), Inches(0.4),
            "CATE estimates from CausalForestDML are highly stable across independent data splits.",
            font_size=14, color=RGBColor(0xCC, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# PART 7: RECOMMENDATIONS
# ═══════════════════════════════════════════════════════════════════════════════
section_header_slide("Part 7", "Strategic Recommendations", "26")

# Slide: City-specific playbook
slide = content_slide("City-Specific Promotional Playbook", "27")

city_recs = [
    ("Mumbai", "Highest Response", [
        "Primary: Coupons (224) + Display (163)",
        "Synergy: SMS + Display during Ganesh Chaturthi & Diwali",
        "Avoid: Deep discounts >20% (revenue destructive)",
    ], ACCENT_ORANGE),
    ("Bengaluru", "Tech-Savvy", [
        "Primary: Coupons (219) + Display (164)",
        "Digital: Consider app-based targeting over SMS",
        "Festival: Pongal week for max seasonal amplification",
    ], ACCENT_BLUE),
    ("Delhi", "Northern Hub", [
        "Primary: Coupons (215) + Display (165)",
        "Synergy: SMS + Display during Diwali & Holi",
        "Cost play: Local Channel Promo highest CATE here (85)",
    ], ACCENT_GREEN),
    ("Hyderabad", "Emerging Market", [
        "Primary: Coupons (211) + SMS Blast (96)",
        "Festival: Eid & Bonalu for maximum impact",
        "Growth: Loyalty programs show above-avg response",
    ], RGBColor(0x9B, 0x59, 0xB6)),
]

for i, (city, tag, bullets, color) in enumerate(city_recs):
    x = Inches(0.4) + i * Inches(3.2)
    card = add_shape(slide, x, Inches(1.3), Inches(3.0), Inches(5.5),
                     RGBColor(0x14, 0x14, 0x25), border_color=color, border_width=Pt(2))

    # City name
    add_rect(slide, x, Inches(1.3), Inches(3.0), Inches(0.8), color)
    add_textbox(slide, x, Inches(1.35), Inches(3.0), Inches(0.45),
                city, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(1.75), Inches(3.0), Inches(0.3),
                tag, font_size=11, color=RGBColor(0xDD, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

    add_bullet_textbox(slide, x + Inches(0.15), Inches(2.4), Inches(2.7), Inches(4.0),
                       bullets, font_size=11, color=RGBColor(0xCC, 0xCC, 0xDD),
                       line_spacing=1.6, bullet_char=">")


# Slide: Do and Don't
slide = content_slide("Action Items: Do's and Don'ts", "28")

dos = [
    "Pair SMS + In-Store Display in all cities (+90 unit synergy)",
    "Cap discount depth at 15% (saturation beyond this point)",
    "Increase coupon distribution in Mumbai & Delhi (highest CATE)",
    "Activate 2x Points during festival weeks (cost-effective amplifier)",
    "Shift budget from discounts to In-Store Display (volume + revenue)",
    "Build city-specific promotional calendars",
]

donts = [
    "Run deep discounts (>20%) for revenue targets -- they destroy value",
    "Run standalone SMS without Display pairing -- 3x less effective",
    "Apply uniform promotions across cities -- CATE varies 10-15%",
    "Omit brand equity from targeting models -- causes 35% bias",
]

# Do column
add_shape(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_GREEN, border_width=Pt(2))
add_textbox(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.5),
            "DO", font_size=24, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
                   dos, font_size=13, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.5,
                   bullet_char="+")

# Don't column
add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), RGBColor(0x14, 0x14, 0x25),
          border_color=ACCENT_RED, border_width=Pt(2))
add_textbox(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.5),
            "DON'T", font_size=24, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(4.5),
                   donts, font_size=13, color=RGBColor(0xCC, 0xCC, 0xDD), line_spacing=1.8,
                   bullet_char="X")


# ═══════════════════════════════════════════════════════════════════════════════
# FINAL: Thank You / Q&A
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.1), ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), ACCENT_BLUE)

add_textbox(slide, Inches(0), Inches(2.2), SLIDE_W, Inches(1.0),
            "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
            font_name="Calibri Light")
add_textbox(slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.8),
            "Questions & Discussion", font_size=28, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER,
            font_name="Calibri Light")

# Tech stack footer
add_multiline_textbox(slide, Inches(2.5), Inches(5.0), Inches(8.0), Inches(1.5), [
    "Tech Stack:  DoWhy  |  EconML  |  scikit-learn  |  pandas  |  matplotlib",
    "Methods:  DAGs  |  Backdoor Criterion  |  Double Machine Learning  |  CausalForest",
    "Repo:  scripts/ + notebooks/ + reports/  |  Fully reproducible (seed=42)",
], font_size=13, color=GREY, line_spacing=1.5, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Presentation saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
